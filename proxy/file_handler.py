"""
File handler for LinkPower - extract text/content from uploaded files
Supports: txt, docx, pptx, pdf, images
"""
import base64
import io
import os

TEXT_EXTS = {'.txt', '.md', '.csv', '.json', '.xml', '.yaml', '.yml', '.log', '.py', '.js', '.ts', '.html', '.css'}
DOC_EXTS = {'.docx', '.pptx', '.pdf'}
IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp'}
MAX_FILE_SIZE = 15 * 1024 * 1024

def get_file_type(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext in TEXT_EXTS: return 'text'
    elif ext in DOC_EXTS: return 'document'
    elif ext in IMAGE_EXTS: return 'image'
    return 'unsupported'

def extract_text_from_txt(data: bytes) -> str:
    for enc in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
        try: return data.decode(enc)
        except UnicodeDecodeError: continue
    return data.decode('utf-8', errors='replace')

def extract_text_from_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip(): lines.append(cell.text.strip())
    return '\n'.join(lines)

def extract_text_from_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip(): lines.append(para.text.strip())
    return '\n'.join(lines)

def extract_text_from_pdf(data: bytes) -> str:
    try:
        import fitz
        doc = fitz.open(stream=data, filetype='pdf')
        parts = [page.get_text().strip() for page in doc if page.get_text().strip()]
        doc.close()
        return '\n'.join(parts)
    except ImportError:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        parts = [page.extract_text().strip() for page in reader.pages if page.extract_text() and page.extract_text().strip()]
        return '\n'.join(parts)

def process_image(data: bytes, filename: str):
    ext = os.path.splitext(filename)[1].lower()
    mm = {'.jpg':'image/jpeg','.jpeg':'image/jpeg','.png':'image/png','.gif':'image/gif','.webp':'image/webp','.bmp':'image/bmp'}
    mime = mm.get(ext, 'image/png')
    return mime, base64.b64encode(data).decode('ascii')

def process_upload(data: bytes, filename: str):
    if len(data) > MAX_FILE_SIZE:
        return {"error": f"文件太大，最大支持 {MAX_FILE_SIZE//1024//1024}MB"}
    ft = get_file_type(filename)
    if ft == 'text':
        text = extract_text_from_txt(data)
        if len(text) > 50000: text = text[:50000] + '\n\n[...内容被截断，最多50000字]'
        return {'type':'text', 'content':text, 'filename':filename, 'size':len(data)}
    if ft == 'document':
        ext = os.path.splitext(filename)[1].lower()
        try:
            if ext == '.docx': text = extract_text_from_docx(data)
            elif ext == '.pptx': text = extract_text_from_pptx(data)
            elif ext == '.pdf': text = extract_text_from_pdf(data)
            else: return {'error': f'不支持的文件格式: {ext}'}
        except Exception as e:
            return {'error': f'文件解析失败: {str(e)}'}
        if len(text) > 50000: text = text[:50000] + '\n\n[...内容被截断，最多50000字]'
        return {'type':'text', 'content':text, 'filename':filename, 'size':len(data)}
    if ft == 'image':
        mime, b64 = process_image(data, filename)
        return {'type':'image', 'mime':mime, 'content':b64, 'filename':filename, 'size':len(data)}
    return {'error': f'不支持的文件类型: {filename}'}
